import io
import logging
import os
import difflib
from pathlib import Path
from typing import Any

import docx
import pdfplumber
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text(filename: str, data: bytes) -> str:
    """Extract text from supported document types based on file extension."""
    ext = Path(filename).suffix.lower()
    if ext == ".txt":
        return data.decode("utf-8")
    if ext == ".docx":
        document = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in document.paragraphs)
    if ext == ".pdf":
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    if ext == ".pptx":
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    raise ValueError(f"Unsupported file type: {ext}")

def diff_texts(old: str, new: str) -> str:
    """Return a unified diff between two strings."""
    diff = difflib.unified_diff(
        old.splitlines(), new.splitlines(), fromfile="old", tofile="new", lineterm=""
    )
    return "\n".join(diff)

def _parse_key_entry(entry: str) -> tuple[str, str]:
    entry = entry.strip()
    if not entry:
        return "", ""
    if ":" in entry:
        provider, key = entry.split(":", 1)
        return provider.lower().strip(), key.strip()
    return "openai", entry


def _gather_api_keys() -> list[tuple[str, str]]:
    parsed: list[tuple[str, str]] = []
    env_keys = os.getenv("SUMMARIZER_API_KEYS", "")
    if env_keys:
        for entry in env_keys.split(","):
            provider, key = _parse_key_entry(entry)
            if provider and key:
                parsed.append((provider, key))
    legacy_key = os.getenv("OPENAI_API_KEY")
    if legacy_key:
        parsed.append(("openai", legacy_key))
    gemini_key = os.getenv("GEMINI_API_KEY")
    if gemini_key:
        parsed.append(("gemini", gemini_key))
    return parsed


def _create_openai_client(api_key: str):
    """Return an OpenAI client instance or None if the SDK is unavailable."""
    try:
        from openai import OpenAI

        return OpenAI(api_key=api_key)
    except ModuleNotFoundError:
        logger.warning("openai SDK is not installed, skipping OpenAI summarizer")
    except Exception as exc:
        logger.warning("failed to create OpenAI client: %s", exc)
    return None


def _fallback_summary(diff_text: str) -> str:
    return "Summary:\n" + diff_text[:500]


def _merge_api_keys(
    override: None | str | list[str]
) -> list[tuple[str, str]]:
    if isinstance(override, list):
        parsed: list[tuple[str, str]] = []
        for entry in override:
            provider, key = _parse_key_entry(entry)
            if provider and key:
                parsed.append((provider, key))
        if parsed:
            return parsed
    if isinstance(override, str) and override.strip():
        provider, key = _parse_key_entry(override)
        if provider and key:
            return [(provider, key)]
    return _gather_api_keys()


def _call_openai(prompt: str, key: str, client_factory=None) -> tuple[str | None, dict[str, Any]]:
    client_factory = client_factory or _create_openai_client
    client = client_factory(key)
    if client is None:
        return None, {}
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=200,
        )
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("OpenAI summarization failed: %s", exc)
        return None, {}
    usage = getattr(response, "usage", None)
    tokens = getattr(usage, "total_tokens", None) if usage else None
    return response.choices[0].message.content.strip(), {"tokens_used": tokens}


def _gemini_model_name() -> str:
    return os.getenv("GEMINI_MODEL", "gemini-2.0-flash-exp")


def _call_gemini(prompt: str, key: str) -> tuple[str | None, dict[str, Any]]:
    try:
        import google.generativeai as genai
    except ModuleNotFoundError:
        logger.warning("google.generativeai SDK not installed, skipping Gemini summarizer")
        return None, {}
    try:
        genai.configure(api_key=key)
        model = genai.GenerativeModel(_gemini_model_name())
        response = model.generate_content(prompt)
        text = getattr(response, "text", None)
        if not text and hasattr(response, "candidates"):
            for candidate in response.candidates:
                parts = getattr(candidate, "content", {}).parts if hasattr(candidate, "content") else []
                if parts:
                    text = getattr(parts[0], "text", None)
                    if text:
                        break
        return text, {}
    except Exception as exc:
        logger.warning("Gemini summarization failed: %s", exc)
        return None, {}


def summarize_changes(
    diff_text: str,
    mission_context: str | None = None,
    client_factory=None,
    api_keys_override: None | str | list[str] = None,
) -> tuple[str, dict[str, Any]]:
    """Summarize diff text using OpenAI/Gemini if available, otherwise return fallback."""
    api_keys = _merge_api_keys(api_keys_override)
    prompt_lines = [
        "You are a medical science liaison translating clinical and promotional updates "
        "into clear, actionable insights for marketing, medical affairs, legal, and sales."
    ]
    if mission_context:
        prompt_lines.append(f"Mission context: {mission_context}")
    prompt_lines.append(
        "Describe changes in a neutral, friendly tone highlighting compliance impact, patient safety, and sales alignment."
    )
    prompt_lines.append("Changes:")
    prompt_lines.append(diff_text)
    client_factory = client_factory or _create_openai_client
    prompt = "\n\n".join(prompt_lines)

    metadata: dict[str, Any] = {
        "method": "fallback",
        "tokens_used": None,
        "truncated": len(diff_text) > 500,
    }

    if not api_keys:
        logger.info("no summarizer API keys configured, using fallback text")
        return _fallback_summary(diff_text), metadata

    for provider, key in api_keys:
        summary = None
        extra_metadata: dict[str, Any] = {}
        if provider == "openai":
            summary, extra_metadata = _call_openai(prompt, key, client_factory)
        elif provider == "gemini":
            summary, extra_metadata = _call_gemini(prompt, key)
        if summary:
            metadata["method"] = provider
            metadata["tokens_used"] = extra_metadata.get("tokens_used")
            metadata["truncated"] = False
            return summary, metadata
        logger.warning("summarization failed with provider %s", provider)

    return _fallback_summary(diff_text), metadata
